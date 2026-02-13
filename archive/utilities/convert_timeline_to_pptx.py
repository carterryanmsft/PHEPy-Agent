"""
Convert Incident Timeline Markdown to PowerPoint Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def create_section_header(prs, title):
    """Create section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section header layout
    slide.shapes.title.text = title
    return slide

def create_content_slide(prs, title, content_lines):
    """Create content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for line in content_lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14)
    
    return slide

def create_metrics_slide(prs, title, metrics):
    """Create metrics dashboard slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add metrics boxes
    box_width = Inches(1.6)
    box_height = Inches(1.2)
    top_margin = Inches(1.5)
    left_margin = Inches(0.8)
    spacing = Inches(0.3)
    
    for idx, (metric_name, metric_value) in enumerate(metrics.items()):
        left = left_margin + (idx % 5) * (box_width + spacing)
        top = top_margin + (idx // 5) * (box_height + spacing)
        
        # Create box
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 120, 215)  # Microsoft blue
        shape.line.color.rgb = RGBColor(0, 90, 158)
        
        # Add text
        text_frame = shape.text_frame
        text_frame.clear()
        
        # Metric name
        p1 = text_frame.paragraphs[0]
        p1.text = metric_name
        p1.font.size = Pt(12)
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER
        
        # Metric value
        p2 = text_frame.add_paragraph()
        p2.text = metric_value
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def create_table_slide(prs, title, headers, rows):
    """Create slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Add table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(0.3) * rows_count
    
    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    
    # Set column widths
    col_width = int(width / cols_count)
    for col_idx in range(cols_count):
        table.columns[col_idx].width = col_width
    
    # Fill header
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 120, 215)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Fill data rows
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    return slide

def parse_markdown_timeline(md_file):
    """Parse markdown file and extract content"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Extract key sections
    sections = {}
    
    # Extract title
    title_match = re.search(r'^# (.+)$', content, re.MULTILINE)
    sections['title'] = title_match.group(1) if title_match else "Incident Timeline"
    
    # Extract incident info
    incident_match = re.search(r'\*\*Incident:\*\* (.+)$', content, re.MULTILINE)
    duration_match = re.search(r'\*\*Duration:\*\* (.+)$', content, re.MULTILINE)
    customer_match = re.search(r'\*\*Customer:\*\* (.+)$', content, re.MULTILINE)
    
    sections['incident'] = incident_match.group(1) if incident_match else ""
    sections['duration'] = duration_match.group(1) if duration_match else ""
    sections['customer'] = customer_match.group(1) if customer_match else ""
    
    # Extract metrics from dashboard table
    metrics = {}
    dashboard_section = re.search(r'## 📊 Executive Metrics Dashboard(.+?)---', content, re.DOTALL)
    if dashboard_section:
        table_match = re.search(r'\|(.+?)\|(.+?)\|(.+?)\|(.+?)\|(.+?)\|', dashboard_section.group(1))
        if table_match:
            headers = [h.strip() for h in table_match.groups()]
            values_match = re.search(r'\|\s*\*\*(.+?)\*\*\s*\|\s*\*\*(.+?)\*\*\s*\|\s*\*\*(.+?)\*\*\s*\|\s*\*\*(.+?)\*\*\s*\|\s*\*\*(.+?)\*\*\s*\|', dashboard_section.group(1))
            if values_match:
                values = [v.strip() for v in values_match.groups()]
                metrics = dict(zip(headers, values))
    
    sections['metrics'] = metrics
    
    return sections

def create_presentation(md_file, output_file):
    """Create PowerPoint from markdown timeline"""
    print(f"📄 Reading markdown file: {md_file}")
    sections = parse_markdown_timeline(md_file)
    
    print(f"🎨 Creating PowerPoint presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    print("  ✓ Creating title slide")
    subtitle = f"{sections['incident']}\n{sections['customer']} | {sections['duration']}"
    create_title_slide(prs, sections['title'], subtitle)
    
    # Slide 2: Executive Metrics Dashboard
    if sections['metrics']:
        print("  ✓ Creating metrics dashboard")
        create_metrics_slide(prs, "📊 Executive Metrics Dashboard", sections['metrics'])
    
    # Slide 3: Section Header - Timeline Overview
    print("  ✓ Creating timeline overview")
    create_section_header(prs, "📅 Timeline Overview")
    
    # Slide 4: Phase Breakdown
    print("  ✓ Creating phase breakdown")
    phases = [
        "Phase 1: Initial Investigation (Dec 13-22) - 9 days",
        "Phase 2: Holiday Communication Gap (Dec 23-29) - 7 days",
        "Phase 3: Customer Pushback (Dec 30 - Jan 6) - 8 days",
        "Phase 4: Multi-ICM Pursuit (Jan 7-16) - 10 days",
        "Phase 5: Deep Dive & Discovery (Jan 17-27) - 11 days",
        "Phase 6: Resolution Path (Jan 28-30) - 3 days"
    ]
    create_content_slide(prs, "Timeline Phases", phases)
    
    # Slide 5: Communication Gaps
    print("  ✓ Creating communication gaps slide")
    gaps = [
        "Gap 1 (Dec 23-29): 7-day holiday communication gap",
        "Gap 2 (Jan 7): Contradictory support messages",
        "Gap 3 (Jan 7): ICM downgraded without clear ownership",
        "Gap 4 (Jan 9-11): No PG response at Sev 2 for 1.48 days",
        "Gap 5 (Jan 12): Classification team engaged too late",
        "Gap 6 (Jan 13-17): No PG response at Sev 2 for 4 days",
        "Gap 7 (Jan 18): Incorrect tenant SIT analysis provided"
    ]
    create_content_slide(prs, "🚨 Communication Gaps Identified", gaps)
    
    # Slide 6: Key Metrics
    print("  ✓ Creating key metrics slide")
    key_metrics = {
        "Total Duration": "48 days",
        "ICM Creation": "3 days",
        "Sev 25 Escalation": "4.7 days",
        "Root Cause": "33 days",
        "Breakthrough": "45 days"
    }
    create_metrics_slide(prs, "⏱️ Key Timeline Metrics", key_metrics)
    
    # Slide 7: Waste Analysis
    print("  ✓ Creating waste analysis")
    waste_data = [
        "Total Duration: 48 days",
        "Value-Add Time: 25 days (52%)",
        "Waste Time: 23 days (48%)",
        "",
        "Major Delays:",
        "• 14 days waiting on PG (multiple gaps)",
        "• 7 days holiday communication gap",
        "• 3 days to create ICM",
        "• 10 days late Classification engagement"
    ]
    create_content_slide(prs, "💡 Waste Analysis", waste_data)
    
    # Slide 8: Critical Events
    print("  ✓ Creating critical events timeline")
    headers = ["Date", "Event", "Impact"]
    rows = [
        ["Dec 13", "Case created", "Incident start"],
        ["Dec 16", "ICM-1 created (3-day delay)", "Escalation delay"],
        ["Dec 22", "Classification timeout identified", "Initial diagnosis"],
        ["Dec 23-29", "Holiday gap - 7 days", "Communication failure"],
        ["Jan 15", "Customer escalates (CFL)", "Executive visibility"],
        ["Jan 16", "ICM-3 Classification team", "Late specialist engagement"],
        ["Jan 27", "100K exception discovered", "Breakthrough"],
        ["Jan 30", "Case downgraded to Sev B", "Resolution"]
    ]
    create_table_slide(prs, "🎯 Critical Path Events", headers, rows)
    
    # Slide 9: Lessons Learned - What Went Wrong
    print("  ✓ Creating lessons learned")
    wrong = [
        "Process Failures:",
        "• 3-day delay creating ICM",
        "• 41-hour delay raising severity",
        "• No engineering bridge for CFL",
        "",
        "Communication Failures:",
        "• 7-day holiday gap with no technical updates",
        "• Multiple PG silence periods at Sev 2",
        "",
        "Technical Failures:",
        "• Classification team engaged 10 days late",
        "• Hidden 100K exception not discovered until Day 45"
    ]
    create_content_slide(prs, "❌ What Went Wrong", wrong)
    
    # Slide 10: Lessons Learned - What Went Right
    print("  ✓ Creating success factors")
    right = [
        "Team Collaboration:",
        "• CAT engagement critical in breakthrough discovery",
        "• Customer persistence drove accountability",
        "",
        "Technical Excellence:",
        "• 10-day telemetry analysis identified patterns",
        "• Classification PG confirmed inefficient regex",
        "• Exception removed within hours of discovery",
        "",
        "Process Improvements:",
        "• DCR raised for UI visibility improvements",
        "• TSG updates identified for Support training"
    ]
    create_content_slide(prs, "✅ What Went Right", right)
    
    # Slide 11: Preventive Measures
    print("  ✓ Creating preventive measures")
    prevention = [
        "Immediate Actions (0-30 days):",
        "• Automated SLA violation alerts",
        "• Holiday coverage protocol with mandatory updates",
        "• Auto-create engineering bridges for CFL",
        "• Exception visibility enhancement",
        "",
        "Short-Term (30-90 days):",
        "• Early specialist engagement criteria",
        "• Communication quality checks",
        "• Telemetry-driven escalation",
        "",
        "Long-Term (90+ days):",
        "• Proactive SIT optimization program",
        "• Secure By Default GA",
        "• Support training & TSG updates"
    ]
    create_content_slide(prs, "🛡️ Preventive Measures", prevention)
    
    # Slide 12: Breakthrough Discovery
    print("  ✓ Creating breakthrough slide")
    breakthrough = [
        "🎯 The Hidden Root Cause:",
        "",
        "Legacy 100K unique SIT exception discovered on Jan 27",
        "",
        "• Pre-500 SIT limit exception persisted",
        "• Kept classification running until 100K unique SIT counts",
        "• Not visible in standard telemetry",
        "• Discovered through deep CAT investigation",
        "",
        "Immediate Action:",
        "• Customer LT approved removal",
        "• PG removed exception at 5 PM ET same day",
        "• Telemetry improvements observed immediately"
    ]
    create_content_slide(prs, "🚀 Breakthrough Discovery - Day 45", breakthrough)
    
    # Slide 13: Summary
    print("  ✓ Creating summary slide")
    summary = [
        "48-day incident with 3 ICMs and 7 communication gaps",
        "",
        "Root Cause: Hidden 100K unique SIT exception",
        "",
        "Resolution: Exception removed + Secure By Default preview",
        "",
        "Key Findings:",
        "• 48% waste (23 days) vs 52% value-add (25 days)",
        "• 14 days waiting on PG across multiple gaps",
        "• Late Classification team engagement (Day 33)",
        "• CAT breakthrough discovery on Day 45",
        "",
        "Customer Status: Sev B - optimizing with CAT support"
    ]
    create_content_slide(prs, "📋 Summary", summary)
    
    # Save presentation
    print(f"💾 Saving presentation: {output_file}")
    prs.save(output_file)
    print(f"✅ PowerPoint created successfully!")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    md_file = r"c:\Users\carterryan\OneDrive - Microsoft\PHEPy\INCIDENT_TIMELINE_2512120040008759_Desjardins_Full.md"
    output_file = r"c:\Users\carterryan\OneDrive - Microsoft\PHEPy\Desjardins_Incident_Timeline.pptx"
    
    create_presentation(md_file, output_file)
